import pptx
from pptx.util import Inches
from lxml import etree
import json

def remove_bullet_point(placeholders):
    for ph in placeholders:
        for paragraph in ph.text_frame.paragraphs:
            paragraph.font.bold = False
            paragraph.font.size = pptx.util.Pt(35)
            paragraph._pPr.insert(0, etree.Element("{http://schemas.openxmlformats.org/drawingml/2006/main}buNone"))

def generate_pptx(songs, output_file):
    prs = pptx.Presentation()

    for song in songs:
        slide_layout = prs.slide_layouts[5]  # Utilise le layout pour un titre et un contenu

        # Ajoute un slide avec le titre de la chanson
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = song['title']

        # Add background image to the slide
        img_path = "background_image.png"  # Specify the path to your background image
        left = top = Inches(0)
        pic = slide.shapes.add_picture(img_path, left, top, width=prs.slide_width, height=prs.slide_height)
        slide.shapes._spTree.remove(pic._element)
        slide.shapes._spTree.insert(2, pic._element)

        # Ajoute un slide pour chaque couplet de la chanson
        for verse in song['verses']:
            slide_layout = prs.slide_layouts[1]  # Utilise un layout de titre et de contenu pour les couplets
            slide = prs.slides.add_slide(slide_layout)
            placeholders = slide.placeholders
            if placeholders[1].text == "":
                content = placeholders[1]
            else:
                content = placeholders[0]
            content.text = '\n'.join(verse['lines'])
            remove_bullet_point([content])
            
            # Add number and title of the song on top of each verse
            top_content = placeholders[0]
            top_content.text = f"{song['number']}. {song['title']}"

            # Add background image to the verse slide
            left = top = Inches(0)
            pic = slide.shapes.add_picture(img_path, left, top, width=prs.slide_width, height=prs.slide_height)
            slide.shapes._spTree.remove(pic._element)
            slide.shapes._spTree.insert(2, pic._element)

    prs.save(output_file)

def main():
    input_file = 'liste_chants.json'
    output_file = 'chants.pptx'

    with open(input_file, 'r', encoding='utf-8') as file:
        songs = json.load(file)

    generate_pptx(songs, output_file)

if __name__ == "__main__":
    main()
